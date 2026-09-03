"""GuruFlow teacher-brain API.

Runs the whole product loop:
``Understand -> Plan -> Explain -> Demonstrate -> Question -> Evaluate ->
Adapt -> Continue``

The API also serves the web client and Person 3's media/visual ES modules, so
the entire SaaS runs from one process with no build step and no API keys.
"""

from __future__ import annotations

import os
import sys
import tempfile
from pathlib import Path
from typing import Any

REPO_ROOT = Path(__file__).resolve().parents[2]
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))

# Load optional .env (GEMINI_API_KEY, TTS/avatar keys) before any service
# that reads os.environ at import time. Secrets are never committed.
try:
    from dotenv import load_dotenv
    load_dotenv(Path(__file__).resolve().parent / ".env")
    load_dotenv(REPO_ROOT / ".env")
except Exception:
    pass

from fastapi import FastAPI, HTTPException, Request, UploadFile, File  # noqa: E402
from fastapi.middleware.cors import CORSMiddleware  # noqa: E402
from fastapi.responses import FileResponse, Response  # noqa: E402
from fastapi.staticfiles import StaticFiles  # noqa: E402

from apps.api.models import (  # noqa: E402
    AnswerRequest,
    EvaluationResult,
    LearningReport,
    LessonPlan,
    MaterialRequest,
    PlanRequest,
)
from apps.api.store import InMemoryLessonRepository, LessonSession  # noqa: E402
from apps.api.student_memory import StudentMemoryStore  # noqa: E402
from services.evaluation import (  # noqa: E402
    CHECKPOINT_CONCEPT_ID,
    CONSTANT_CURRENT,
    DIRECT_PROPORTIONALITY,
    build_repair_scene,
    build_report,
    evaluate_answer,
)
from services.ingestion import ingest_text, ingest_topic, Material  # noqa: E402
from services.planner import plan_lesson  # noqa: E402
from services.planner.flashcards import generate_flashcards  # noqa: E402
from services.planner.persona import persona_feedback  # noqa: E402
from services.planner.study_plan import build_study_plan  # noqa: E402
from services.qa import answer_question  # noqa: E402
from services import video as video_service  # noqa: E402

WEB_DIR = REPO_ROOT / "apps" / "web"

# Mastery credited for a concept the learner watched all the way through.
TAUGHT_MASTERY = 0.8

app = FastAPI(
    title="GuruFlow Teacher Brain",
    version="2.0.0",
    description=(
        "Source-grounded, multilingual AI teacher. Plans lessons from material, "
        "detects misconceptions and adapts the next scene."
    ),
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=os.getenv("GURUFLOW_CORS_ORIGINS", "*").split(","),
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

repository = InMemoryLessonRepository()

# Long-term, cross-session student memory. Survives server restarts.
student_memory = StudentMemoryStore()


@app.middleware("http")
async def no_store_for_assets(request, call_next):
    """Stop browsers caching the client during a demo.

    A stale styles.css or app.js after a last-minute fix is a classic demo
    failure; the app is small enough that revalidating every asset costs
    nothing.
    """
    response = await call_next(request)
    path = request.url.path
    if path.endswith((".js", ".css", ".html", ".json")) or path == "/":
        response.headers["Cache-Control"] = "no-store, must-revalidate"
    return response


# ---------------------------------------------------------------------------
# Health and materials
# ---------------------------------------------------------------------------


@app.get("/health")
def health() -> dict[str, Any]:
    from services.llm import gemini_available
    from services.rag import rag_status
    return {
        "status": "ok",
        "service": "guruflow-teacher-brain",
        "mode": "deterministic" if not gemini_available() else "llm-enhanced",
        "gemini": gemini_available(),
        "rag": rag_status(),
        "video": video_service.cache_stats(),
    }


@app.post("/materials")
def create_material(request: MaterialRequest) -> dict[str, Any]:
    """Ingest a topic or pasted text and report extraction status."""
    if request.text:
        material = ingest_text(request.text, request.title or "Uploaded material")
    else:
        material = ingest_topic(request.topic or "")
    repository.save_material(material)
    return material.to_dict()


@app.get("/materials/{material_id}")
def get_material(material_id: str) -> dict[str, Any]:
    material = repository.get_material(material_id)
    if material is None:
        raise HTTPException(status_code=404, detail="Material not found")
    return material.to_dict()


# ---------------------------------------------------------------------------
# Lesson planning
# ---------------------------------------------------------------------------


@app.post("/lessons/plan", response_model=LessonPlan)
def create_plan(request: PlanRequest) -> dict[str, Any]:
    """Plan a lesson for this learner, grounded in the chosen material."""
    material = None
    if request.materialId:
        material = repository.get_material(request.materialId)
        if material is None:
            raise HTTPException(status_code=404, detail="Material not found")
    if material is None:
        material = ingest_topic(request.topic)
        repository.save_material(material)

    learner = request.learner.model_dump(exclude_none=True)
    plan = plan_lesson(
        learner, material, topic=request.topic, study_mode=request.studyMode
    )

    session = LessonSession(
        lesson_id=plan["id"],
        student_id=request.studentId,
        plan=plan,
        material_id=material.material_id,
    )
    repository.save_session(session)
    return plan


@app.get("/lessons/{lesson_id}", response_model=LessonPlan)
def get_lesson(lesson_id: str) -> dict[str, Any]:
    session = repository.get_session(lesson_id)
    if session is None:
        raise HTTPException(status_code=404, detail="Lesson not found")
    return session.plan


@app.post("/lessons/{lesson_id}/language", response_model=LessonPlan)
def switch_language(lesson_id: str, body: dict[str, str]) -> dict[str, Any]:
    """Re-render the lesson in another language, keeping learner progress.

    Progress lives on the session, not on the plan, so re-planning in a new
    language never resets mastery, attempts or diagnosed misconceptions.
    """
    session = repository.get_session(lesson_id)
    if session is None:
        raise HTTPException(status_code=404, detail="Lesson not found")

    language = body.get("language")
    if language not in ("english", "hindi", "hinglish"):
        raise HTTPException(status_code=400, detail="Unsupported language")

    material = repository.get_material(session.material_id)
    if material is None:
        raise HTTPException(status_code=404, detail="Material not found")

    learner = dict(session.plan["learner"])
    learner["language"] = language
    session.plan = plan_lesson(
        learner,
        material,
        topic=session.plan.get("topic", "Ohm's Law"),
        lesson_id=session.lesson_id,
    )
    repository.save_session(session)
    return session.plan


@app.post("/lessons/{lesson_id}/scenes/{scene_id}/complete")
def complete_scene(lesson_id: str, scene_id: str) -> dict[str, Any]:
    """Mark a scene watched so the final report reflects real progress."""
    session = repository.get_session(lesson_id)
    if session is None:
        raise HTTPException(status_code=404, detail="Lesson not found")

    scene = next((s for s in session.plan["scenes"] if s["id"] == scene_id), None)

    if scene is None:
        # Repair scenes are generated during evaluation and never stored on the
        # plan, but the learner did watch them, so count them as progress
        # without crediting mastery for a concept they got wrong.
        if scene_id.startswith("scene-repair-"):
            session.scenes_completed.add(scene_id)
            repository.save_session(session)
            return {
                "lessonId": lesson_id,
                "sceneId": scene_id,
                "scenesCompleted": len(session.scenes_completed),
                "totalScenes": len(session.plan["scenes"]),
                "repair": True,
            }
        raise HTTPException(status_code=404, detail="Scene not found")

    session.scenes_completed.add(scene_id)
    concept_id = scene["conceptId"]
    # Mastery is credited for being taught something. A checkpoint is assessed,
    # not taught, so watching one earns nothing -- its concept's mastery is
    # owned by evaluation.
    is_assessment = bool(scene.get("checkpointId")) or concept_id == CHECKPOINT_CONCEPT_ID
    if not is_assessment and concept_id not in session.concept_mastery:
        session.set_mastery(concept_id, TAUGHT_MASTERY)
    repository.save_session(session)

    return {
        "lessonId": lesson_id,
        "sceneId": scene_id,
        "scenesCompleted": len(session.scenes_completed),
        "totalScenes": len(session.plan["scenes"]),
    }


# ---------------------------------------------------------------------------
# Evaluation and adaptation
# ---------------------------------------------------------------------------


@app.post(
    "/lessons/{lesson_id}/checkpoints/{checkpoint_id}/answer",
    response_model=EvaluationResult,
    # Omit null optionals so an "advance" result carries no empty repairScene.
    response_model_exclude_none=True,
)
def answer_checkpoint(
    lesson_id: str, checkpoint_id: str, request: AnswerRequest
) -> dict[str, Any]:
    """Evaluate a checkpoint answer and decide what to teach next."""
    session = repository.get_session(lesson_id)
    if session is None:
        raise HTTPException(status_code=404, detail="Lesson not found")

    language = request.language or session.plan["learner"]["language"]
    attempt = session.record_attempt(checkpoint_id)
    result = evaluate_answer(
        request.answer, language=language, attempt=attempt, option_id=request.optionId
    )

    # Apply the learner's chosen teacher personality to the feedback tone.
    personality = session.plan["learner"].get("personality")
    if personality:
        tone = persona_feedback(
            personality, encouraged=bool(result.get("correct"))
        )
        if tone:
            result["feedback"] = f"{tone}{result['feedback']}"

    session.set_mastery(CHECKPOINT_CONCEPT_ID, result["mastery"])

    if result["nextAction"] == "repair":
        session.checkpoints_failed += 1
        session.note_misconception(result["misconception"], CHECKPOINT_CONCEPT_ID)
    elif result["correct"]:
        session.checkpoints_passed += 1
        # A correct answer after a repair closes the misconception.
        session.resolve_misconceptions()

    repository.save_session(session)

    result["attempt"] = attempt
    return result


# ---------------------------------------------------------------------------
# Reporting
# ---------------------------------------------------------------------------


def _report_for_session(session: LessonSession) -> dict[str, Any]:
    return build_report(
        student_id=session.student_id,
        lesson_id=session.lesson_id,
        concept_mastery=session.concept_mastery,
        misconceptions=list(session.misconceptions.values()),
        scenes_completed=len(session.scenes_completed),
        checkpoints_passed=session.checkpoints_passed,
        checkpoints_failed=session.checkpoints_failed,
        total_time_seconds=session.elapsed_seconds(),
    )


@app.get("/lessons/{lesson_id}/report", response_model=LearningReport)
def lesson_report(lesson_id: str) -> dict[str, Any]:
    session = repository.get_session(lesson_id)
    if session is None:
        raise HTTPException(status_code=404, detail="Lesson not found")
    report = _report_for_session(session)
    # Fold this finished lesson into the student's long-term memory.
    student_memory.record_lesson(session, report)
    return report


@app.get("/students/{student_id}/report", response_model=LearningReport)
def student_report(student_id: str) -> dict[str, Any]:
    sessions = repository.sessions_for_student(student_id)
    if not sessions:
        raise HTTPException(status_code=404, detail="No lessons for this student")
    report = _report_for_session(sessions[-1])
    # Fold into long-term memory so the profile reflects the latest lesson.
    student_memory.record_lesson(sessions[-1], report)
    return report


# ---------------------------------------------------------------------------
# Long-term student memory and review features
# ---------------------------------------------------------------------------


@app.get("/students/{student_id}/profile")
def student_profile(student_id: str) -> dict[str, Any]:
    """Return a student's long-term learning profile across all lessons.

    Survives server restarts (file-backed). Includes running concept mastery,
    recurring misconceptions and lesson history for the dashboard.
    """
    profile = student_memory.get_profile(student_id)
    if profile is None:
        raise HTTPException(status_code=404, detail="No learning profile yet")
    return profile


@app.get("/students/{student_id}/history")
def student_history(student_id: str) -> dict[str, Any]:
    """Return the chronological lesson history for a student."""
    profile = student_memory.get_profile(student_id)
    if profile is None:
        raise HTTPException(status_code=404, detail="No learning history yet")
    return {"studentId": student_id, "lessons": profile.get("lessons", [])}


@app.post("/lessons/{lesson_id}/flashcards")
def lesson_flashcards(lesson_id: str, body: dict[str, Any] | None = None) -> dict[str, Any]:
    """Generate review flashcards from a lesson's scenes.

    Optional body may select a ``language`` and/or filter ``conceptIds`` (e.g.
    the learner's weak concepts from long-term memory). The core Ohm's Law
    formula card is always included so a review is complete.
    """
    session = repository.get_session(lesson_id)
    if session is None:
        raise HTTPException(status_code=404, detail="Lesson not found")

    body = body or {}
    language = body.get("language") or session.plan["learner"]["language"]
    concept_ids = body.get("conceptIds")

    cards = generate_flashcards(
        session.plan["scenes"],
        language=language,
        concept_ids=concept_ids,
    )
    return {
        "lessonId": lesson_id,
        "studentId": session.student_id,
        "language": language,
        "count": len(cards),
        "cards": cards,
    }


@app.get("/students/{student_id}/study-plan")
def student_study_plan(student_id: str) -> dict[str, Any]:
    """Return a spaced, multi-day revision plan built from long-term memory.

    Uses the student's accumulated concept mastery and recurring weak concepts
    to schedule review sessions across one week (days 1, 2, 4, 7), revising
    weak concepts earliest and most often. 404 if the student has no memory
    yet (i.e. they need to finish at least one lesson first).
    """
    profile = student_memory.get_profile(student_id)
    if profile is None:
        raise HTTPException(status_code=404, detail="No learning profile yet")
    return build_study_plan(
        student_id,
        concept_mastery=profile.get("conceptMastery", {}),
        weak_concepts=profile.get("weakConcepts", []),
    )


# ---------------------------------------------------------------------------
# Text-to-Speech
# ---------------------------------------------------------------------------

VOICE_MAP = {
    "english": "en-IN-MadhurNeural",
    "hindi": "hi-IN-SwaraNeural",
    "hinglish": "hi-IN-MadhurNeural",
}


@app.post("/tts")
async def text_to_speech(body: dict[str, str]) -> Response:
    """Generate speech audio from text using edge-tts (free, no API key)."""
    text = body.get("text", "").strip()
    language = body.get("language", "hinglish")
    if not text:
        raise HTTPException(status_code=400, detail="Text is required")

    voice = VOICE_MAP.get(language, VOICE_MAP["hinglish"])

    try:
        import edge_tts
        communicate = edge_tts.Communicate(text, voice)
        tmp = tempfile.NamedTemporaryFile(suffix=".mp3", delete=False)
        tmp_path = tmp.name
        tmp.close()
        await communicate.save(tmp_path)
        return FileResponse(tmp_path, media_type="audio/mpeg", filename="speech.mp3")
    except ImportError:
        raise HTTPException(status_code=503, detail="edge-tts not installed. Run: pip install edge-tts")
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"TTS generation failed: {str(exc)}")


# ---------------------------------------------------------------------------
# File upload for PDF parsing
# ---------------------------------------------------------------------------


@app.post("/upload")
async def upload_material(file: UploadFile = File(...)) -> dict[str, Any]:
    """Upload a PDF, DOCX, PPTX or text file and extract sections.

    Sections are the contract-shaped ``material.sections`` used by retrieval,
    so every supported format feeds the same citation/grounding pipeline.
    """
    content = await file.read()
    filename = file.filename or "uploaded-file"
    lower_name = filename.lower()

    sections: list[dict[str, Any]] | None = None
    if lower_name.endswith(".pdf"):
        sections = _parse_pdf(content)
    elif lower_name.endswith(".docx"):
        sections = _parse_docx(content)
    elif lower_name.endswith((".ppt", ".pptx")):
        sections = _parse_pptx(content)
    elif lower_name.endswith((".txt", ".md")):
        text = content.decode("utf-8", errors="replace")
        material = ingest_text(text, title=filename)
        repository.save_material(material)
        return material.to_dict()
    else:
        raise HTTPException(
            status_code=400,
            detail="Supported formats: PDF, DOCX, PPTX, TXT, MD",
        )

    # Persist the parsed document as a real, retrievable Material so the plan
    # endpoint can look it up by materialId (must match what we return).
    material = Material(
        material_id=f"material-upload-{filename}",
        document_id=f"material-upload-{filename}",
        title=filename,
        status="ready",
        sections=sections,
        origin="upload",
    )
    repository.save_material(material)
    return material.to_dict()


@app.post("/diagram")
async def render_diagram(body: dict[str, Any]) -> Response:
    """Render a scene's visual spec as a real PNG diagram.

    Accepts ``{"type": "circuit"|"equation"|"graph"|"concept_map"|...}`` plus
    optional ``title``/``nodes`` data. Returns image/png bytes so the visual a
    scene calls for is genuinely rendered, not just described.
    """
    from apps.api.diagram import render_diagram as _render

    try:
        png = _render(body, size=(640, 360))
    except Exception as exc:  # renderer is defensive, but never let a 500 leak raw
        raise HTTPException(status_code=500, detail=f"Diagram render failed: {exc}")
    return Response(
        content=png,
        media_type="image/png",
        headers={},
    )


def _parse_pdf(content: bytes) -> list[dict[str, Any]]:
    """Extract one section per non-empty PDF page using PyMuPDF."""
    try:
        import fitz
    except ImportError:
        raise HTTPException(
            status_code=503,
            detail="PyMuPDF not installed. Run: pip install PyMuPDF",
        )
    doc = fitz.open(stream=content, filetype="pdf")
    sections: list[dict[str, Any]] = []
    try:
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            if text.strip():
                lines = text.strip().splitlines()
                heading = lines[0].strip()[:80] if lines else f"Page {page_num + 1}"
                sections.append(_make_section(
                    index=page_num,
                    page_or_slide=page_num + 1,
                    heading=heading,
                    text=text,
                ))
    finally:
        doc.close()
    return sections


def _parse_docx(content: bytes) -> list[dict[str, Any]]:
    """Extract document sections from a Word .docx file using python-docx.

    Splits on heading-style paragraphs so structural headings become distinct
    sections with page numbers approximated by paragraph index (docs have no
    hard pagination, so we number by block order instead).
    """
    try:
        from io import BytesIO
        from docx import Document
    except ImportError:
        raise HTTPException(
            status_code=503,
            detail="python-docx not installed. Run: pip install python-docx",
        )
    document = Document(BytesIO(content))
    sections: list[dict[str, Any]] = []
    current_heading: str | None = None
    current_blocks: list[str] = []

    def flush() -> None:
        nonlocal current_heading, current_blocks
        text = "\n".join(current_blocks).strip()
        if text:
            sections.append(_make_section(
                index=len(sections),
                page_or_slide=len(sections) + 1,
                heading=current_heading or f"Section {len(sections) + 1}",
                text=text,
            ))
        current_heading = None
        current_blocks = []

    for para in document.paragraphs:
        style = (para.style.name or "").lower()
        text = para.text.strip()
        if not text:
            continue
        if style.startswith("heading") or style == "title":
            flush()
            current_heading = text
        else:
            current_blocks.append(text)
    flush()
    return sections


def _parse_pptx(content: bytes) -> list[dict[str, Any]]:
    """Extract one section per slide from a PowerPoint .pptx file."""
    try:
        from io import BytesIO
        from pptx import Presentation
    except ImportError:
        raise HTTPException(
            status_code=503,
            detail="python-pptx not installed. Run: pip install python-pptx",
        )
    prs = Presentation(BytesIO(content))
    sections: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(prs.slides):
        blocks: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        blocks.append(text)
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        blocks.append(" | ".join(cells))
        text = "\n".join(blocks).strip()
        if text:
            heading = blocks[0][:80] if blocks else f"Slide {slide_index + 1}"
            sections.append(_make_section(
                index=slide_index,
                page_or_slide=slide_index + 1,
                heading=heading,
                text=text,
            ))
    return sections


def _make_section(
    index: int,
    page_or_slide: int,
    heading: str,
    text: str,
) -> dict[str, Any]:
    """Build a contract-shaped section with derived keywords."""
    return {
        "sectionId": f"upload-sec-{index + 1}",
        "pageOrSlide": page_or_slide,
        "heading": heading,
        "excerpt": text[:500],
        "keywords": _derive_upload_keywords(text),
    }


_STOPWORDS_UPLOAD = {
    "the", "a", "an", "of", "to", "and", "is", "are", "in", "on", "for",
    "it", "that", "this", "with", "as", "by", "be", "or", "from", "at",
}

def _derive_upload_keywords(text: str) -> list[str]:
    import re
    words = re.findall(r"[a-zA-Z]{3,}", text.lower())
    seen: list[str] = []
    for word in words:
        if word not in _STOPWORDS_UPLOAD and word not in seen:
            seen.append(word)
    return seen[:12]


# ---------------------------------------------------------------------------
# Static hosting: web client plus Person 3's ES modules
# ---------------------------------------------------------------------------

# Serving the media/visuals sources lets the browser import the real contract
# modules instead of a reimplementation of them.
# ---------------------------------------------------------------------------
# Follow-up questions
# ---------------------------------------------------------------------------


@app.post("/lessons/{lesson_id}/ask")
def ask_followup(lesson_id: str, body: dict[str, Any]) -> dict[str, Any]:
    """Answer a learner's follow-up without losing the lesson's context.

    Grounded in the same material the lesson was planned from, so the answer
    carries page citations - or admits the material does not cover it.
    """
    session = repository.get_session(lesson_id)
    if session is None:
        raise HTTPException(status_code=404, detail="Lesson not found")

    question = str(body.get("question", "")).strip()
    if not question:
        raise HTTPException(status_code=400, detail="A question is required")

    material = repository.get_material(session.material_id)
    if material is None:
        raise HTTPException(status_code=404, detail="Material not found")

    language = body.get("language") or session.plan["learner"]["language"]
    result = answer_question(
        question,
        material.sections,
        material.document_id,
        language=language,
        lesson_topic=session.plan.get("topic", "this lesson"),
    )
    result["question"] = question
    return result


# ---------------------------------------------------------------------------
# Teaching video
# ---------------------------------------------------------------------------


def _video_payload(video_id: str) -> dict[str, Any]:
    status = video_service.get_status(video_id)
    return {
        "videoId": video_id,
        "status": status,
        "url": f"/video/{video_id}.mp4" if status == "ready" else None,
    }


@app.post("/video/render")
def request_scene_video(body: dict[str, Any]) -> dict[str, Any]:
    """Ask for a scene's teaching video, rendering in the background.

    Takes the Scene itself rather than an id so repair scenes - which are
    generated during evaluation and never belong to the stored plan - get
    videos on exactly the same path as planned scenes.
    """
    scene = body.get("scene")
    if not isinstance(scene, dict) or not scene.get("narration"):
        raise HTTPException(status_code=400, detail="A scene object is required")

    language = body.get("language") or "hinglish"
    if not video_service.video_generation_available():
        raise HTTPException(
            status_code=503,
            detail="Video generation unavailable (manim or ffmpeg missing)",
        )

    return _video_payload(video_service.render_in_background(scene, language))


@app.get("/video/{video_id}/status")
def scene_video_status(video_id: str) -> dict[str, Any]:
    return _video_payload(video_id)


@app.get("/video/{video_id}.mp4")
def serve_scene_video(video_id: str, request: Request) -> Response:
    """Serve a rendered video, honouring Range so the player can seek."""
    path = video_service.cached_path(video_id)
    if path is None:
        raise HTTPException(status_code=404, detail="Video not ready")

    size = path.stat().st_size
    range_header = request.headers.get("range")
    if not range_header:
        return FileResponse(path, media_type="video/mp4")

    # "bytes=START-END"; END is optional.
    try:
        units, _, span = range_header.partition("=")
        if units.strip().lower() != "bytes":
            raise ValueError(units)
        start_text, _, end_text = span.partition("-")
        start = int(start_text) if start_text else 0
        end = int(end_text) if end_text else size - 1
    except ValueError:
        raise HTTPException(status_code=416, detail="Malformed Range header")

    start = max(0, start)
    end = min(end, size - 1)
    if start > end:
        raise HTTPException(status_code=416, detail="Range not satisfiable")

    with path.open("rb") as handle:
        handle.seek(start)
        chunk = handle.read(end - start + 1)

    return Response(
        content=chunk,
        status_code=206,
        media_type="video/mp4",
        headers={
            "Content-Range": f"bytes {start}-{end}/{size}",
            "Accept-Ranges": "bytes",
            "Content-Length": str(len(chunk)),
        },
    )


@app.post("/lessons/{lesson_id}/video/prerender")
def prerender_lesson_videos(lesson_id: str) -> dict[str, Any]:
    """Warm the whole lesson's video cache so the demo never waits."""
    session = repository.get_session(lesson_id)
    if session is None:
        raise HTTPException(status_code=404, detail="Lesson not found")
    if not video_service.video_generation_available():
        raise HTTPException(status_code=503, detail="Video generation unavailable")

    language = session.plan["learner"]["language"]
    scenes = list(session.plan["scenes"])

    # The repair scene is the flagship moment, so warm it too rather than
    # making the learner wait for it at the worst possible time.
    for misconception in (DIRECT_PROPORTIONALITY, CONSTANT_CURRENT):
        scenes.append(build_repair_scene(misconception, language))

    ids = video_service.prerender_lesson(scenes, language)
    return {"lessonId": lesson_id, "videoIds": ids, "count": len(ids)}


@app.get("/lessons/{lesson_id}/video/status")
def lesson_video_status(lesson_id: str) -> dict[str, Any]:
    session = repository.get_session(lesson_id)
    if session is None:
        raise HTTPException(status_code=404, detail="Lesson not found")
    language = session.plan["learner"]["language"]
    scenes = session.plan["scenes"]
    statuses = [
        {"sceneId": scene["id"], **_video_payload(video_service.video_id(scene, language))}
        for scene in scenes
    ]
    ready = sum(1 for s in statuses if s["status"] == "ready")
    return {
        "lessonId": lesson_id,
        "ready": ready,
        "total": len(statuses),
        "scenes": statuses,
    }


app.mount(
    "/vendor/visuals",
    StaticFiles(directory=REPO_ROOT / "services" / "visuals" / "src"),
    name="visuals",
)
app.mount(
    "/vendor/media",
    StaticFiles(directory=REPO_ROOT / "services" / "media" / "src"),
    name="media",
)
app.mount(
    "/fixtures",
    StaticFiles(directory=REPO_ROOT / "demo-fixtures"),
    name="fixtures",
)


_FAVICON_SVG = (
    b'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 32 32">'
    b'<rect width="32" height="32" rx="7" fill="#6c8cff"/>'
    b'<text x="16" y="22" font-family="sans-serif" font-size="14" '
    b'font-weight="bold" fill="white" text-anchor="middle">GF</text></svg>'
)


@app.get("/favicon.ico")
def favicon() -> Response:
    """Browsers request /favicon.ico even when a data-URI icon is declared.

    Serving it turns a red 404 in the console into a non-event, which matters
    when a judge has devtools open.
    """
    return Response(content=_FAVICON_SVG, media_type="image/svg+xml")


@app.get("/")
def index() -> FileResponse:
    return FileResponse(WEB_DIR / "index.html")


app.mount("/", StaticFiles(directory=WEB_DIR, html=True), name="web")
