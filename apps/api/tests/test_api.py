"""End-to-end API tests for the judge-critical path.

Upload/topic -> plan -> scene -> wrong answer -> detected misconception ->
repair -> correct retry -> final report.
"""

from __future__ import annotations

import sys
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[3]
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))

import pytest  # noqa: E402
from fastapi.testclient import TestClient  # noqa: E402

from apps.api.main import app, repository  # noqa: E402


@pytest.fixture(autouse=True)
def deterministic_engine(monkeypatch):
    """Force the deterministic planner for these endpoint tests.

    These assert the deterministic engine's exact structure (e.g. scene count
    == 7, page-numbered citations). With a live GEMINI_API_KEY set, the LLM
    path would return a variable number of scenes and make the assertions
    flaky. Suppress the key and reset the cached Gemini client so every run is
    hermetic, exactly as test_planner.py does.
    """
    import services.llm as llm

    saved_client = llm._gemini_client
    saved_attempted = llm._model_attempted
    monkeypatch.delenv("GEMINI_API_KEY", raising=False)
    monkeypatch.delenv("GURUFLOW_LLM_API_KEY", raising=False)
    llm._gemini_client = None
    llm._model_attempted = False
    try:
        yield
    finally:
        llm._gemini_client = saved_client
        llm._model_attempted = saved_attempted


def make_plan(client, **learner_overrides):
    learner = {
        "level": "beginner",
        "language": "hinglish",
        "availableMinutes": 20,
        "goal": "Understand Ohm's Law",
    }
    learner.update(learner_overrides)
    response = client.post(
        "/lessons/plan", json={"learner": learner, "topic": "Ohm's Law"}
    )
    assert response.status_code == 200, response.text
    return response.json()


def _import_fails(module_name: str) -> bool:
    try:
        __import__(module_name)
        return False
    except ImportError:
        return True


def test_health(client):
    response = client.get("/health")
    assert response.status_code == 200
    assert response.json()["status"] == "ok"


def test_material_ingestion_reports_pages(client):
    response = client.post("/materials", json={"topic": "Electricity"})
    assert response.status_code == 200
    body = response.json()
    assert body["status"] == "ready"
    assert body["sectionCount"] > 0
    assert body["pageCount"] > 0
    assert body["documentId"] == "ncert-class9-science-ch12"


def test_plan_endpoint_returns_contract_shaped_lesson(client):
    plan = make_plan(client)
    assert plan["id"].startswith("lesson-")
    assert plan["learner"]["language"] == "hinglish"
    assert len(plan["scenes"]) == 7
    assert plan["scenes"][0]["citations"][0]["pageOrSlide"] >= 1


def test_lesson_can_be_fetched_again(client):
    plan = make_plan(client)
    fetched = client.get(f"/lessons/{plan['id']}")
    assert fetched.status_code == 200
    assert fetched.json()["id"] == plan["id"]


def _study_mode_plan(client, mode):
    learner = {
        "level": "beginner",
        "language": "hinglish",
        "availableMinutes": 20,
        "goal": "Understand Ohm's Law",
    }
    response = client.post(
        "/lessons/plan",
        json={"learner": learner, "topic": "Ohm's Law", "studyMode": mode},
    )
    assert response.status_code == 200, response.text
    return response.json()


def test_exam_study_mode_produces_drilled_assessment_plan(client):
    plan = _study_mode_plan(client, "exam")
    assert plan["studyMode"] == "exam"
    ids = [s["conceptId"] for s in plan["scenes"]]
    # Assessment-heavy: extra practice drill, the real checkpoint and review.
    assert "ohms-law-practice" in ids
    assert "ohms-law-application" in ids
    assert "lesson-summary" in ids
    assert any(s.get("checkpointId") for s in plan["scenes"])


def test_revision_study_mode_produces_quick_recap(client):
    plan = _study_mode_plan(client, "revision")
    assert plan["studyMode"] == "revision"
    ids = [s["conceptId"] for s in plan["scenes"]]
    assert "ohms-law" in ids
    assert "ohms-law-application" in ids
    # Revision skips the deep lecture-only scenes.
    assert "intro-electricity" not in ids


def test_default_lesson_mode_still_returns_time_tier(client):
    plan = _study_mode_plan(client, "lesson")
    assert plan["studyMode"] == "lesson"
    assert len(plan["scenes"]) == 7  # full tier at 20 minutes


def test_study_plan_endpoint_requires_profile_then_builds_schedule(client):
    # No memory yet -> 404 with a helpful message.
    assert client.get("/students/student-demo/study-plan").status_code == 404
    # Finish a lesson to seed long-term memory.
    plan = make_plan(client)
    lesson_id = plan["id"]
    for scene in plan["scenes"]:
        client.post(f"/lessons/{lesson_id}/scenes/{scene['id']}/complete")
    client.get(f"/lessons/{lesson_id}/report")
    # Now a spaced revision plan is available.
    res = client.get("/students/student-demo/study-plan")
    assert res.status_code == 200
    body = res.json()
    assert body["strategy"] == "spaced-repetition"
    assert body["horizonDays"] == 7
    assert len(body["sessions"]) == 4
    assert all(s["mode"] == "revision" for s in body["sessions"])


def test_judge_critical_path_wrong_then_repair_then_retry(client):
    plan = make_plan(client)
    lesson_id = plan["id"]
    checkpoint_id = next(s["checkpointId"] for s in plan["scenes"] if s.get("checkpointId"))

    # Watch every teaching scene before the checkpoint.
    for scene in plan["scenes"]:
        client.post(f"/lessons/{lesson_id}/scenes/{scene['id']}/complete")

    # 1. Deliberately wrong answer -> misconception detected -> repair scene.
    wrong = client.post(
        f"/lessons/{lesson_id}/checkpoints/{checkpoint_id}/answer",
        json={"answer": "Current increases when resistance increases."},
    ).json()
    assert wrong["correct"] is False
    assert wrong["misconception"] == "direct-proportionality confusion"
    assert wrong["nextAction"] == "repair"
    assert wrong["repairScene"]["visual"]["data"]["analogy"] == "water-pipe"

    # 2. Correct retry -> advance.
    retry = client.post(
        f"/lessons/{lesson_id}/checkpoints/{checkpoint_id}/answer",
        json={"answer": "current kam hoga"},
    ).json()
    assert retry["correct"] is True
    assert retry["nextAction"] == "advance"
    assert retry["attempt"] == 2

    # 3. Final report reflects the whole journey.
    report = client.get(f"/lessons/{lesson_id}/report").json()
    assert report["checkpointsFailed"] == 1
    assert report["checkpointsPassed"] == 1
    assert report["scenesCompleted"] == 7
    assert report["misconceptions"][0]["id"] == "direct-proportionality confusion"
    assert report["misconceptions"][0]["status"] == "resolved"
    assert report["nextTopic"]["title"] == "Series and Parallel Circuits"


def test_answer_changes_the_next_scene(client):
    """Guards against 'fake adaptation': the answer must alter what comes next."""
    plan = make_plan(client)
    lesson_id = plan["id"]
    checkpoint_id = next(s["checkpointId"] for s in plan["scenes"] if s.get("checkpointId"))

    correct = client.post(
        f"/lessons/{lesson_id}/checkpoints/{checkpoint_id}/answer",
        json={"answer": "Current decreases"},
    ).json()

    plan2 = make_plan(client)
    lesson2 = plan2["id"]
    checkpoint2 = next(s["checkpointId"] for s in plan2["scenes"] if s.get("checkpointId"))
    wrong = client.post(
        f"/lessons/{lesson2}/checkpoints/{checkpoint2}/answer",
        json={"answer": "Current increases"},
    ).json()

    assert correct["nextAction"] != wrong["nextAction"]
    assert "repairScene" not in correct
    assert "repairScene" in wrong


def test_mcq_answer_path(client):
    plan = make_plan(client)
    lesson_id = plan["id"]
    checkpoint_id = next(s["checkpointId"] for s in plan["scenes"] if s.get("checkpointId"))

    result = client.post(
        f"/lessons/{lesson_id}/checkpoints/{checkpoint_id}/answer",
        json={"answer": "", "optionId": "increases"},
    ).json()
    assert result["nextAction"] == "repair"


def test_language_switch_preserves_progress(client):
    plan = make_plan(client)
    lesson_id = plan["id"]
    checkpoint_id = next(s["checkpointId"] for s in plan["scenes"] if s.get("checkpointId"))

    client.post(f"/lessons/{lesson_id}/scenes/{plan['scenes'][0]['id']}/complete")
    client.post(
        f"/lessons/{lesson_id}/checkpoints/{checkpoint_id}/answer",
        json={"answer": "Current increases"},
    )

    switched = client.post(
        f"/lessons/{lesson_id}/language", json={"language": "english"}
    ).json()

    assert switched["id"] == lesson_id
    assert switched["learner"]["language"] == "english"
    assert "Hello students" in switched["scenes"][0]["narration"]

    # Progress survived the switch.
    report = client.get(f"/lessons/{lesson_id}/report").json()
    assert report["scenesCompleted"] == 1
    assert report["checkpointsFailed"] == 1
    assert report["misconceptions"][0]["status"] == "open"


@pytest.mark.parametrize("minutes", [5, 10, 20, 60])
def test_time_budget_matrix_through_api(client, minutes):
    plan = make_plan(client, availableMinutes=minutes)
    assert plan["estimatedSeconds"] <= minutes * 60
    assert any(s.get("checkpointId") for s in plan["scenes"])


@pytest.mark.parametrize("level", ["beginner", "intermediate", "advanced"])
@pytest.mark.parametrize("language", ["english", "hindi", "hinglish"])
def test_level_language_matrix_through_api(client, level, language):
    plan = make_plan(client, level=level, language=language)
    assert plan["learner"]["level"] == level
    assert plan["learner"]["language"] == language
    for scene in plan["scenes"]:
        assert scene["narration"].strip()
        assert scene["objective"].strip()


def test_watching_a_checkpoint_earns_no_mastery(client):
    """A checkpoint is assessed, not taught, so watching it proves nothing."""
    plan = make_plan(client)
    lesson_id = plan["id"]
    checkpoint_scene = next(s for s in plan["scenes"] if s.get("checkpointId"))

    client.post(f"/lessons/{lesson_id}/scenes/{checkpoint_scene['id']}/complete")
    report = client.get(f"/lessons/{lesson_id}/report").json()

    assert checkpoint_scene["conceptId"] not in report["strongConcepts"]
    assert checkpoint_scene["conceptId"] not in report["weakConcepts"]


def test_repair_scene_counts_as_watched(client):
    """Repair scenes are not on the plan but the learner still watched them."""
    plan = make_plan(client)
    lesson_id = plan["id"]
    checkpoint_id = next(s["checkpointId"] for s in plan["scenes"] if s.get("checkpointId"))

    wrong = client.post(
        f"/lessons/{lesson_id}/checkpoints/{checkpoint_id}/answer",
        json={"answer": "Current increases when resistance increases."},
    ).json()
    repair_id = wrong["repairScene"]["id"]

    response = client.post(f"/lessons/{lesson_id}/scenes/{repair_id}/complete")
    assert response.status_code == 200
    assert response.json()["repair"] is True

    # An unknown, non-repair scene id is still a 404.
    assert client.post(f"/lessons/{lesson_id}/scenes/not-a-scene/complete").status_code == 404


def test_missing_lesson_returns_404(client):
    assert client.get("/lessons/does-not-exist").status_code == 404
    assert (
        client.post(
            "/lessons/nope/checkpoints/cp/answer", json={"answer": "x"}
        ).status_code
        == 404
    )


def test_student_report_endpoint(client):
    plan = make_plan(client)
    client.post(f"/lessons/{plan['id']}/scenes/{plan['scenes'][0]['id']}/complete")
    report = client.get("/students/student-demo/report")
    assert report.status_code == 200
    assert report.json()["studentId"] == "student-demo"


@pytest.mark.skipif(
    _import_fails("docx"), reason="python-docx not installed"
)
def test_docx_upload_parses_sections(client):
    from docx import Document
    from io import BytesIO

    doc = Document()
    doc.add_heading("Resistance Fundamentals", level=1)
    doc.add_paragraph(
        "Resistance opposes the flow of current and is measured in ohms."
    )
    doc.add_heading("Ohm's Law", level=1)
    doc.add_paragraph("Ohm's law says V equals I times R.")
    buf = BytesIO()
    doc.save(buf)

    res = client.post(
        "/upload",
        files={"file": ("material.docx", buf.getvalue(),
                        "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
    )
    assert res.status_code == 200, res.text
    body = res.json()
    assert body["sectionCount"] == 2
    assert body["sections"][0]["heading"] == "Resistance Fundamentals"
    assert body["sections"][0]["pageOrSlide"] == 1
    assert body["sections"][0]["sectionId"] == "upload-sec-1"
    assert "resistance" in body["sections"][0]["keywords"]


@pytest.mark.skipif(
    _import_fails("pptx"), reason="python-pptx not installed"
)
def test_pptx_upload_parses_sections(client):
    from pptx import Presentation
    from io import BytesIO

    prs = Presentation()
    st = prs.slides.add_slide(prs.slide_layouts[1])
    st.shapes.title.text = "Current"
    st.placeholders[1].text = "Current is the flow of charge in amperes."

    buf = BytesIO()
    prs.save(buf)

    res = client.post(
        "/upload",
        files={"file": ("deck.pptx", buf.getvalue(),
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    assert res.status_code == 200, res.text
    body = res.json()
    assert body["sectionCount"] == 1
    assert body["sections"][0]["pageOrSlide"] == 1
    assert body["sections"][0]["heading"] == "Current"


def test_unsupported_upload_format_rejected(client):
    res = client.post(
        "/upload",
        files={"file": ("notes.xlsx", b"not a supported doc", "application/octet-stream")},
    )
    assert res.status_code == 400


def test_uploaded_text_material_can_be_planned_end_to_end(client):
    """Regression test for the upload -> plan chain (Bug 1).

    An uploaded file must be persisted so the materialId returned by /upload
    is accepted by /lessons/plan. Previously the parsed sections were returned
    to the browser but never saved server-side, so planning 404'd.
    """
    raw = (
        "Ohm's Law\n"
        "The current through a conductor is directly proportional to the "
        "voltage across it when resistance is constant.\n\n"
        "Series Circuits\n"
        "In a series circuit, the same current flows through every component.\n"
    )
    up = client.post("/upload", files={"file": ("notes.txt", raw.encode(), "text/plain")})
    assert up.status_code == 200, up.text
    material_id = up.json()["materialId"]
    assert material_id

    plan_resp = client.post(
        "/lessons/plan",
        json={
            "learner": {
                "level": "beginner", "language": "english", "availableMinutes": 10,
                "goal": "Understand Ohm's Law",
            },
            "topic": "Ohm's Law",
            "materialId": material_id,
        },
    )
    assert plan_resp.status_code == 200, plan_resp.text
    plan = plan_resp.json()
    assert plan["materialId"] == material_id
    # The uploaded material grounds at least one scene with a real citation.
    assert any(scene["citations"] for scene in plan["scenes"])


def test_uploaded_docx_material_can_be_planned_end_to_end(client):
    """The parse-and-persist path must also work for binary formats (Bug 1)."""
    from io import BytesIO
    from docx import Document

    doc = Document()
    doc.add_heading("resistance", level=1)
    doc.add_paragraph("Resistance is the property of a conductor to resist the flow of charges through it.")
    buf = BytesIO()
    doc.save(buf)

    up = client.post(
        "/upload",
        files={"file": ("notes.docx", buf.getvalue(),
                        "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
    )
    assert up.status_code == 200, up.text
    material_id = up.json()["materialId"]

    plan_resp = client.post(
        "/lessons/plan",
        json={
            "learner": {
                "level": "beginner", "language": "english", "availableMinutes": 10,
                "goal": "Understand resistance",
            },
            "topic": "Resistance",
            "materialId": material_id,
        },
    )
    assert plan_resp.status_code == 200, plan_resp.text
    assert plan_resp.json()["materialId"] == material_id


def test_unknown_topic_returns_honest_refusal_not_electricity(client):
    """Regression test for Bug 2: asking for an unsupported topic without an
    LLM key must never silently teach a different subject (Electricity)."""
    response = client.post(
        "/lessons/plan",
        json={
            "learner": {
                "level": "beginner", "language": "hinglish", "availableMinutes": 10,
                "goal": "Learn it",
            },
            "topic": "Photosynthesis",
        },
    )
    assert response.status_code == 200, response.text
    plan = response.json()
    assert plan["unsupportedTopic"] is True
    assert plan["tier"] == "unsupported-topic"
    assert all(scene["conceptId"] == "unsupported-topic" for scene in plan["scenes"])
    assert all(scene["citations"] == [] for scene in plan["scenes"])
    # Never teach electricity content under a Photosynthesis label.
    assert all("Electric Current" not in s["narration"] for s in plan["scenes"])


def test_diagram_renders_real_png_for_each_visual_type(client):
    for vtype in ("circuit", "equation", "graph", "concept_map", "water_pipe_analogy"):
        res = client.post("/diagram", json={"type": vtype, "title": vtype})
        assert res.status_code == 200, res.text
        assert res.headers["content-type"].startswith("image/png")
        assert res.content[:4] == b"\x89PNG"
        assert len(res.content) > 1000
